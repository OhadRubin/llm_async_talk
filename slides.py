from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_agent1_slides():
    """Create slides 1-6 of the LLM Async Talk presentation"""
    
    # Create presentation object
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Formalizing Latency in Distributed Systems and Its Role in Multi‑Agent LLM Coordination"
    subtitle.text = "Duration: ~20 minutes\nAudience: Systems researchers, ML engineers, distributed systems practitioners"
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Good morning/afternoon. Today I want to share a journey that started with a simple question: How do we let multiple LLMs have a conversation? This led me down a rabbit hole connecting 40 years of distributed systems theory with the cutting-edge challenges of orchestrating language models. By the end, you'll see why the pharynx - yes, your throat - might be the best mental model for understanding LLM coordination."
    
    # Slide 2: The Problem in One Image (Hook → Stakes)
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "The Problem in One Image"
    tf = body.text_frame
    tf.text = "Uncoordinated LLMs break down."
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "[ANIMATION: 3-Phase Visual Story]"
    p.level = 1
    p.font.italic = True
    p.font.color.rgb = RGBColor(128, 128, 128)
    
    p = tf.add_paragraph()
    p.text = "Phase 1: Three LLMs ready and initialized"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 2: All start generating simultaneously" 
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Phase 3: Text streams collide → Garbled output"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "If we can't coordinate them, ensemble reasoning fails."
    p.level = 0
    p.font.bold = True
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Each LLM must complete its entire response before processing new input"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Architecturally incapable of listening while speaking"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Uncoordinated LLMs break down. [SHOW 3-PHASE ANIMATION] Phase 1: Three LLMs - Claude, GPT-4, and Gemini - are initialized and ready, like processes waiting for a critical section. Phase 2: Without coordination, they all start generating text simultaneously - multiple processes entering the critical section at once. Phase 3: Their text streams collide and create garbled output - this is our race condition resulting in corrupted shared state. If we can't coordinate them, ensemble reasoning fails - multiple models can't vote, debate can't improve accuracy, and simulation breaks down. Unlike humans who can interrupt and adjust, each LLM, once started, must complete its entire response. The visual shows this as a classic coordination problem - like network packet collisions or database transaction conflicts."
    
    # Slide 3: Agenda
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Agenda"
    tf = body.text_frame
    tf.text = "1. Background - Formalizing latency & LLM architecture"
    
    p = tf.add_paragraph()
    p.text = "2. Motivation - Why coordination matters"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "3. Method - The talking stick protocol"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "4. Results - What happened when we tried it"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "5. Implications - Lessons for system design"
    p.level = 0
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "We'll start with the theoretical foundations - how computer science has formalized time and delay. Then we'll see why LLMs break our usual assumptions. I'll show you a protocol we implemented, share some surprising results, and discuss what this means for building multi-agent AI systems."
    
    # Slide 4: Latency Models & Synchrony (Compressed Background)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Latency Models & Synchrony"
    tf = body.text_frame
    tf.text = "Four formal approaches to reasoning about delay:"
    
    p = tf.add_paragraph()
    p.text = "Temporal Logic (requirements) • Network Calculus (bounds) • Process Calculi (composition) • Queuing Theory (statistics)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "The Synchrony Spectrum:"
    p.level = 0
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Asynchronous ←→ Partially Synchronous ←→ Synchronous"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Internet/Email ←→ Datacenter/Raft ←→ CPU bus/Clock"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Key insight: What you can build depends on timing assumptions"
    p.level = 0
    p.font.bold = True
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Computer science offers multiple lenses for reasoning about delay - from temporal logic requirements to statistical queuing models. The key insight from distributed systems: what you can build depends on timing assumptions. Messages arrive 'eventually' on the Internet, 'usually bounded' in datacenters, or 'always within Δ' on CPU buses. Bounded time changes what's possible. Given that, how does it shape LLMs?"
    
    # Slide 5: Why Timing Bounds Matter (Examples)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Why Timing Bounds Matter"
    tf = body.text_frame
    tf.text = "Without bounds (async) vs With bounds (sync):"
    
    p = tf.add_paragraph()
    p.text = "❌ Consensus impossible with 1 failure (FLP theorem)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Byzantine consensus tolerates 1/3 failures"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "❌ Cannot detect failures perfectly"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Timeout = failure (if no heartbeat for 2Δ)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Real example: Credit cards work globally because financial networks have timing bounds"
    p.level = 0
    p.font.bold = True
    p.font.italic = True
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "This isn't academic trivia. The famous FLP theorem proves consensus is impossible in asynchronous systems with even one failure. But add timing bounds? You can tolerate up to 1/3 Byzantine failures and detect failures perfectly. This is why your credit card works at any ATM globally - there are timing bounds in the financial network. Bounded time changes what's possible."
    
    # Slide 6: The Biological Analogy (Moved here for A-B-A' flow)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "The Biological Analogy"
    tf = body.text_frame
    tf.text = "HUMAN THROAT                    LLM PIPELINE"
    
    p = tf.add_paragraph()
    p.text = "     ↓                               ↓"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "[Air/Food] → Pharynx → [Lungs/Stomach]    [Input] → Attention → [Output]"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "              ↑                                        ↑"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "        Single channel                          Single context"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "        Can't do both                          Can't do both"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "So far we saw how timing bounds determine what's possible in distributed systems; now let's see how this connects to LLMs. Here's a mental model that clicked for me. Humans face a similar constraint - we breathe and eat through the same tube. The pharynx is a single point of failure. Evolution's solution? The epiglottis - a valve that switches between modes. We literally hold our breath to swallow. LLMs need a similar mechanism - a protocol to switch between listening and speaking modes."
    
    
    
    
    
    # Slide 7: LLM Architecture Primer
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "LLM Architecture: Two Distinct Phases"
    tf = body.text_frame
    tf.text = "Input: \"What is the capital of France?\""
    
    p = tf.add_paragraph()
    p.text = "         ↓"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "┌─────────────────┐"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "│     PREFILL     │ (Process all input tokens in parallel)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "│  Compute KV     │ Time: ~50ms for 1K tokens"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "│  Cache for all  │ Bottleneck: FLOPS - can be parallelized"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "└─────────────────┘"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "         ↓"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "┌─────────────────┐"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "│    GENERATE     │ (Produce one token at a time)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "│  The... →       │ Time: ~30ms per token"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "│  capital... →   │ Bottleneck: Memory bandwidth"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "│  is... →        │ Each token depends on ALL previous tokens"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "│  Paris.         │ ❌ CANNOT BE INTERRUPTED"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "└─────────────────┘"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "So far we saw the biological analogy for single-channel constraints; now let's see the technical details of how LLMs actually work. LLMs have two distinct phases. Prefill processes your entire prompt in parallel - it's compute-bound, doing massive matrix multiplications. Generate produces one token at a time, each depending on all previous tokens - it's memory-bound, constantly loading cached attention values. Here's the critical constraint: once generation starts, the model cannot process new input until it's done. It's architecturally serial. This is crucial for distributed systems researchers to understand: LLMs are fundamentally different from traditional processes. In prefill, they process your entire prompt in parallel - like a massive matrix multiplication computing attention weights for all input tokens simultaneously. But generation is inherently sequential. Each new token depends on the attention-weighted combination of ALL previous tokens. This isn't a software choice - it's baked into the transformer mathematics. You cannot interrupt generation mid-stream because the attention mechanism requires the complete sequence. This breaks traditional distributed systems assumptions about interruptible processes."
    
    # Slide 8: The Sequential Bottleneck
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "The Sequential Bottleneck"
    tf = body.text_frame
    tf.text = "Time →"
    
    p = tf.add_paragraph()
    p.text = "LLM-A: [PREFILL] → [GENERATING RESPONSE............] → [DONE]"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "                            ↑"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "LLM-B:            \"Hey, wait I want to say—\""
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "                   ❌ Cannot process until generation completes"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Why this breaks distributed systems assumptions:"
    p.level = 0
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Traditional processes: Can be interrupted at any instruction"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Traditional processes: Can receive signals mid-execution"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• LLMs: Each token = f(ALL previous tokens) via attention"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• LLMs: Mathematical dependency chain cannot be broken"
    p.level = 1
    p.font.bold = True
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "This is our core challenge and the key insight for distributed systems researchers. Once an LLM starts generating, it's like a printer from 1995 - you can't add pages to the queue until the current job finishes. Every token depends on all previous tokens through the attention mechanism. There's no architectural way to inject new information mid-stream. This isn't a software limitation or design choice - it's baked into the transformer mathematics. Traditional distributed systems assume processes can be interrupted, can receive signals, can checkpoint state. LLMs violate all these assumptions during generation."
    
    # Slide 9: The Staleness Problem (Moved here for constraint → symptom → motivation flow)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "The Staleness Problem"
    tf = body.text_frame
    tf.text = "T=0s:   A sees: []                    B sees: []"
    
    p = tf.add_paragraph()
    p.text = "T=1s:   A says: \"Hello\"               B generating: \"Hi there...\""
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "T=3s:   A says: \"How are you?\"        B still generating..."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "T=5s:   A says: \"Hello?\"              B completes: \"Hi there, nice to meet you!\""
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "                                      (B never saw A's 2nd and 3rd messages)"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Staleness compounds quickly. LLM-A sends three messages while B generates one response. When B finally speaks, it's responding to ancient history. It's like having a conversation over postal mail - by the time your letter arrives, the context has shifted. In human conversation, we use backchannels - 'uh-huh', 'right' - to stay synchronized. LLMs can't do this. [ENERGY RESET: Ask audience - 'Any guess who grabs the stick first?' to keep engagement at midpoint]"
    
    # Slide 10: Why Coordinate Multiple LLMs?
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Why Coordinate Multiple LLMs?"
    tf = body.text_frame
    tf.text = "1. Ensemble Reasoning: Multiple models vote on answers"
    
    p = tf.add_paragraph()
    p.text = "Like random forests for LLMs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Reduces individual model errors"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Adversarial Debate: Models argue positions to find truth"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Improves factual accuracy through peer review"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Role-Play Simulation: Models embody different perspectives"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Complex social dynamics modeling"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Decision-making simulation"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Why bother? Three compelling use cases. First, ensemble reasoning - like random forests for LLMs, multiple models can vote to reduce errors. Second, adversarial debate - having models argue improves factual accuracy, like peer review. Third, simulation - modeling complex social dynamics or decision-making. But all three break down if responses are based on stale context. Coordination isn't optional."
    
    # Slide 11: Design Space for Coordination
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Design Space: Options We Considered"
    tf = body.text_frame
    tf.text = "Four approaches we explored:"
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Token interrupts (keystroke sync) - maximum responsiveness, terrible efficiency"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Chunk-based turns (walkie-talkie) - speak in paragraphs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Priority requests (speaking queue) - bid for speaking time"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Parallel drafts (operational transform) - merge like Google Docs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "We picked Talking Stick because:"
    p.level = 0
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✓ Simple to implement and understand"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Natural social pressure mechanisms"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✓ Graceful degradation under race conditions"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "So far we saw the problem and motivation for coordination; now let's see our approach. We explored four approaches, each with trade-offs between responsiveness and implementation complexity. Token-level interrupts give maximum responsiveness but terrible efficiency. Chunk-based turns work like a walkie-talkie. Priority requests let agents bid for speaking time. Parallel drafts merge simultaneously like Google Docs. We picked the talking stick because it's simple to implement, has natural social pressure mechanisms, and degrades gracefully under race conditions."
    
    
    
    
    
    
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "The Talking Stick Protocol"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # State diagram section
    diagram_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4))
    diagram_frame = diagram_box.text_frame
    diagram_frame.text = "State Diagram:"
    
    # Add state diagram in monospace font
    p = diagram_frame.add_paragraph()
    p.text = """    ┌─────────────┐
    │   WAITING   │←────────────┐
    └──────┬──────┘             │
           │                    │
      talking_stick()           │
           │                    │
    ┌──────▼──────┐             │
    │ HAS STICK   │             │
    └──────┬──────┘             │
           │                    │
      append(text)              │
           │                    │
    ┌──────▼──────┐             │
    │  COMPOSING  │             │
    └──────┬──────┘             │
           │                    │
        push()                  │
           │                    │
           └────────────────────┘"""
    p.font.name = 'Courier New'
    p.font.size = Pt(12)
    
    # API section
    api_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(4))
    api_frame = api_box.text_frame
    api_frame.text = "API:"
    api_frame.paragraphs[0].font.bold = True
    
    api_methods = [
        "• talking_stick() - Request exclusive speaking rights",
        "• append(text) - Stage message content",
        "• push() - Publish to shared chat", 
        "• check() - Poll for updates"
    ]
    
    for method in api_methods:
        p = api_frame.add_paragraph()
        p.text = method
        p.level = 0
        p.font.size = Pt(14)
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "We implemented the simplest thing that could work - a talking stick. Like the speaking tradition in some Indigenous councils, whoever holds the stick has the floor. The protocol has four primitives. First, request the stick. Then append your message - you can do this multiple times for long messages. Push to publish. And constantly check for updates. It's polling-based, not event-driven, which has implications we'll see. [Note: If time runs short, combine this slide with the next System Feedback slide.]"
    
    # Slide 14: System Feedback Mechanisms
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "System Feedback Mechanisms"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Chat interface mockup
    chat_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(4))
    chat_frame = chat_box.text_frame
    chat_frame.text = "Chat Interface:"
    chat_frame.paragraphs[0].font.bold = True
    chat_frame.paragraphs[0].font.size = Pt(16)
    
    # Chat content with box outline
    p = chat_frame.add_paragraph()
    p.text = """┌─────────────────────────────────────┐
│ Dwight: "Hi everyone!"              │
│                                     │
│ [System: Jim has claimed the        │
│  talking stick and wants to speak]  │
│                                     │
│ [System: Michael has been waiting   │
│  for a response for 4 seconds...]   │
└─────────────────────────────────────┘"""
    p.font.name = 'Courier New'
    p.font.size = Pt(14)
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "The secret sauce is system messages that create social pressure. When someone claims the stick, everyone sees it. When someone's been waiting, there's a gentle nudge. It's like seeing someone's hand raised in a Zoom call - you naturally want to yield. These ambient cues coordinate behavior without strict enforcement."
    
    # Slide 15: The Experiment Setup
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "The Experiment Setup"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Three participant panels
    panels_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    panels_frame = panels_box.text_frame
    panels_frame.text = """┌──────────────┐  ┌──────────────┐  ┌──────────────┐
│   DWIGHT     │  │     JIM      │  │   MICHAEL    │
│              │  │              │  │              │
│ "Interested  │  │ "Philosophy  │  │ "Black holes │
│ in astro-    │  │ and moral    │  │ fascinate    │
│ biology"     │  │ questions"   │  │ me!"         │
└──────────────┘  └──────────────┘  └──────────────┘"""
    panels_frame.paragraphs[0].font.name = 'Courier New'
    panels_frame.paragraphs[0].font.size = Pt(12)
    
    # Task description
    task_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    task_frame = task_box.text_frame
    task_frame.text = "Task: Have a natural conversation"
    task_frame.paragraphs[0].font.size = Pt(16)
    task_frame.paragraphs[0].font.bold = True
    
    p = task_frame.add_paragraph()
    p.text = "Hidden context: One might be a murderer (they don't know this)"
    p.font.size = Pt(14)
    p.font.italic = True
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "So far we saw the method - the talking stick protocol; now let's see what happened when we tested it. We tested with three LLM instances, each with a persona. Dwight's into astrobiology, Jim likes philosophy, Michael loves black holes. We gave them a simple task - have a natural conversation. But here's the twist we didn't tell them - the system was actually trying to identify which one might be a murderer based on conversation patterns. This created an interesting dynamic we'll see unfold."
    
    # Slide 16: Timeline of Key Events
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Timeline of Key Events"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Gantt chart
    gantt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    gantt_frame = gantt_box.text_frame
    gantt_frame.text = "Gantt Chart:"
    gantt_frame.paragraphs[0].font.bold = True
    
    p = gantt_frame.add_paragraph()
    p.text = """Time (s) →  0    2    4    6    8    10   12   14
Dwight:     ■■■──────■■■────────■■■──────■■■
            talk    check      talk     talk

Jim:        ──■■■────────■■■──────■■■────────
             check      talk    check

Michael:    ────■■■──────────■■■──────■■■───
              check         talk    check
            
Events:     └─Intro─┘└Race┘└─Topics─┘└Philosophy┘"""
    p.font.name = 'Courier New'
    p.font.size = Pt(11)
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Here's what actually happened. The introduction phase worked perfectly - clean turn-taking. Then we hit our first race condition at T=6 - both Dwight and Jim claimed the stick simultaneously. The system handled it gracefully by just... letting both messages through. The conversation naturally evolved from interests to deeper philosophical questions about human nature and evil - remember, one might be a murderer."
    
    # Slide 17: Quantitative Results
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Quantitative Results"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Results table
    table_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    table_frame = table_box.text_frame
    table_frame.text = "Dashboard with key metrics:"
    table_frame.paragraphs[0].font.bold = True
    table_frame.paragraphs[0].font.size = Pt(16)
    
    # Table content
    table_data = [
        "| Metric | Value | Context |",
        "|--------|-------|---------|",
        "| Polling frequency | 0.3 Hz | Every ~3 seconds |",
        "| Average response latency | 4.2s | Time holding stick |",
        "| Race conditions | 2 | Both resolved naturally |",
        "| Message truncations | 3 | Due to token limits |",
        "| Conversation coherence | 94% | Human-rated score |",
        "| Total messages | 18 | Over 3 minutes |"
    ]
    
    for row in table_data:
        p = table_frame.add_paragraph()
        p.text = row
        p.font.name = 'Courier New'
        p.font.size = Pt(12)
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "The numbers tell an interesting story. Polling dominated network traffic - every 3 seconds, each LLM checked for updates. The average 'thinking time' while holding the stick was 4.2 seconds. We saw two race conditions where multiple LLMs claimed the stick - both resolved without intervention. Despite truncation issues and races, human raters scored the conversation as 94% coherent. Not bad for a distributed system with no central coordinator!"
    
    # Slide 18: Qualitative Observations
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Qualitative Observations"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Conversation excerpt with annotations
    convo_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    convo_frame = convo_box.text_frame
    convo_frame.text = "Conversation excerpt with annotations:"
    convo_frame.paragraphs[0].font.bold = True
    convo_frame.paragraphs[0].font.size = Pt(16)
    
    # Conversation content
    convo_content = [
        '',
        'Jim: "Have either of you ever encountered real evil?"',
        '     ↑ [Philosophical probe - possibly revealing?]',
        '',
        'Michael: "I think most humans are fundamentally good,',
        '         but circumstances can corrupt."',
        '         ↑ [Optimistic response - deflecting?]',
        '',
        'Dwight: "Environment shapes us... life might exist',
        '        in unimaginable forms"',
        '        ↑ [Bridging science and philosophy]'
    ]
    
    for line in convo_content:
        p = convo_frame.add_paragraph()
        p.text = line
        p.font.name = 'Courier New'
        p.font.size = Pt(12)
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "The conversation took a fascinating turn. Jim steered toward moral philosophy and asked about 'real evil' - significant given the hidden murder context. Michael gave an optimistic response about human nature. Dwight tried to bridge both topics. What emerged wasn't just functional communication but genuinely interesting discourse. The protocol enabled but didn't constrain the natural flow of ideas."
    

    # Slide 19: Emergent Behaviors
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Emergent Behaviors"
    tf = body.text_frame
    tf.text = "1. Conversational Crossed Wires"
    
    p = tf.add_paragraph()
    p.text = "Natural misalignments that happen in real conversation"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Example: Jim answering a question just as it was asked"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Social Pressure Dynamics"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "'Waiting for 4 seconds' created urgency to respond"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Self-regulating without hard limits"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Topic Convergence"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Started with separate interests"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Naturally found common ground in philosophy"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Protocol enabled but didn't force this"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "So far we saw the quantitative and qualitative results from our experiment; now let's see what unexpected behaviors emerged. Three behaviors emerged that we didn't explicitly design for. First, conversational crossed wires - like when Jim revealed his interests just as Dwight asked about them. This felt natural, not like a bug. Second, the waiting notifications created organic pressure to keep things moving. Third, despite starting with different topics - space, philosophy, black holes - the conversation naturally converged. The protocol provided structure without strangling spontaneity. These behaviors inform system design lessons."
    
    # Slide 20: Lessons for System Design
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Lessons for System Design"
    tf = body.text_frame
    tf.text = "Challenge → Traditional Solution → LLM Reality"
    
    p = tf.add_paragraph()
    p.text = "Concurrency → Threads + locks → Sequential by design"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Interruption → Signals/callbacks → Must complete generation"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Coordination → Event-driven → Polling + social cues"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Consistency → ACID/eventual → Embrace crossed wires"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "LLMs break our usual distributed systems playbook. We can't use threads - generation is inherently sequential. We can't interrupt - the attention mechanism requires completing the sequence. Event-driven architectures assume immediate response - LLMs need seconds to think. Perhaps most importantly, perfect consistency might be wrong - those crossed wires made the conversation more human."
    
    # Slide 21: The Staleness-Liveness Tradeoff
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "The Staleness-Liveness Tradeoff"
    tf = body.text_frame
    tf.text = "Fundamental tradeoff in LLM coordination:"
    
    p = tf.add_paragraph()
    p.text = "Liveness ↑"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "    │ Our protocol"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "    │     ★"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "    │   ╱  Ideal (impossible?)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "    │ ╱"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "    └─────────────→ Staleness"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Lower staleness = More interruptions = Less liveness"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Higher liveness = Longer responses = More staleness"
    p.level = 0
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "We face a fundamental tradeoff. Reduce staleness by checking more often, and you sacrifice liveness - conversations become choppy. Improve liveness with longer uninterrupted responses, and staleness accumulates. Our protocol found a middle ground, but the perfect solution might not exist. It's CAP theorem for conversations - you can't have perfect coherence, availability, and partition tolerance."
    
    # Slide 22: Future Directions
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Future Directions"
    tf = body.text_frame
    tf.text = "Research roadmap from current state:"
    
    p = tf.add_paragraph()
    p.text = "Current: Polling + Talking Stick"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "           │"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "           ├─→ Event-Driven Updates"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "           │   (WebSockets, push notifications)"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "           ├─→ Predictive Protocols"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "           │   (Start prefill before your turn)"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "           └─→ Architectural Changes"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "               (Streamable attention?)"
    p.level = 2
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Where do we go from here? Three promising directions. First, replace polling with event-driven updates - reduce that network overhead. Second, predictive protocols - start processing likely responses before it's your turn, like speculative execution for conversation. Most ambitiously, could we modify the transformer architecture itself to allow streaming attention updates? That would solve the root cause, not just manage symptoms."
    
    # Slide 23: Key Takeaways
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Key Takeaways"
    tf = body.text_frame
    tf.text = "1. Latency bounds unlock capabilities"
    
    p = tf.add_paragraph()
    p.text = "From FLP impossibility to Byzantine consensus"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Time assumptions determine what's buildable"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. LLMs have a hard architectural constraint"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Sequential generation is not a bug but a feature"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Need protocols that embrace, not fight this"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Simple protocols can enable complex behavior"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Talking stick + social pressure = coherent conversation"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Emergence happens at the protocol boundary"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. Perfect coordination might be undesirable"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Crossed wires make conversations human"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Some staleness creates authenticity"
    p.level = 1
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Four insights to leave you with. First, the 40-year lesson from distributed systems remains true - time bounds determine what's possible. Second, LLMs' sequential nature isn't a limitation to overcome but a constraint to design around. Third, simple protocols can enable surprisingly complex behavior - our basic talking stick led to philosophical discourse. Finally, and perhaps most surprisingly, perfect coordination might be wrong. Those messy overlaps and misalignments? They're not bugs - they're what make conversation feel real."
    
    # Slide 24: Questions?
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Questions?"
    subtitle.text = "Contact and Resources:\n\n• Protocol implementation (GitHub)\n• Full conversation transcript\n• Related papers on bounded latency\n\nThank you!"
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "Thank you. I'm happy to discuss any aspect - from the formal models to the implementation details to the philosophical questions our LLMs raised about human nature. And if you're wondering - we never did figure out which one was the murderer. Maybe the real mystery was the distributed systems we built along the way."
    
    # Backup Slides
    
    # Slide B1: Formal Staleness Definition
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom content
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Backup Slide B1: Formal Staleness Definition"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "For the theory-inclined:"
    content_frame.paragraphs[0].font.bold = True
    content_frame.paragraphs[0].font.size = Pt(16)
    
    p = content_frame.add_paragraph()
    p.text = ""
    
    p = content_frame.add_paragraph()
    p.text = "Staleness(t) = t - max{t' : t' ≤ t ∧ observed(state(t'))}"
    p.font.name = 'Courier New'
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content_frame.add_paragraph()
    p.text = ""
    
    p = content_frame.add_paragraph()
    p.text = "Where:"
    p.font.size = Pt(14)
    p.font.bold = True
    
    definitions = [
        "• t = current time",
        "• state(t') = global state at time t'",
        "• observed() = when state became visible to agent"
    ]
    
    for definition in definitions:
        p = content_frame.add_paragraph()
        p.text = definition
        p.font.size = Pt(14)
    
    # Slide B2: Protocol Pseudocode
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Backup Slide B2: Protocol Pseudocode"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Code content
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    code_frame = code_box.text_frame
    code_frame.text = """class LLMAgent:
    async def converse(self):
        while not done:
            state = await check()
            if should_speak(state):
                if await talking_stick():
                    msg = generate_response(state)
                    for chunk in tokenize(msg):
                        await append(chunk)
                    await push()
            await sleep(POLL_INTERVAL)"""
    code_frame.paragraphs[0].font.name = 'Courier New'
    code_frame.paragraphs[0].font.size = Pt(12)
    
    # Slide B3: Race Condition Handling
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Backup Slide B3: Race Condition Handling"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "What happens when two agents claim the stick simultaneously?"
    content_frame.paragraphs[0].font.bold = True
    content_frame.paragraphs[0].font.size = Pt(16)
    
    race_steps = [
        "",
        "1. Both receive success (distributed systems are messy)",
        "2. Both compose messages",
        "3. Both push to chat",
        "4. Messages appear in receipt order",
        "5. Conversation continues naturally",
        "",
        "No explicit resolution needed!"
    ]
    
    for step in race_steps:
        p = content_frame.add_paragraph()
        p.text = step
        if step.startswith(("1.", "2.", "3.", "4.", "5.")):
            p.font.size = Pt(14)
        elif step == "No explicit resolution needed!":
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.italic = True
        else:
            p.font.size = Pt(14)
    
    # Slide B4: Transformer Architecture Details (For Technical Deep-Dive)
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Backup Slide B4: Transformer Attention Mechanism"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Technical content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "Mathematical Foundation of Sequential Constraint:"
    content_frame.paragraphs[0].font.bold = True
    content_frame.paragraphs[0].font.size = Pt(16)
    
    p = content_frame.add_paragraph()
    p.text = ""
    
    p = content_frame.add_paragraph()
    p.text = "For token at position i, attention output:"
    p.font.size = Pt(14)
    
    p = content_frame.add_paragraph()
    p.text = "Attention(Q_i, K_{1:i}, V_{1:i}) = softmax(Q_i * K_{1:i}^T / √d_k) * V_{1:i}"
    p.font.name = 'Courier New'
    p.font.size = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = ""
    
    p = content_frame.add_paragraph()
    p.text = "Key Constraints:"
    p.font.bold = True
    p.font.size = Pt(14)
    
    technical_points = [
        "• Causal masking: Token i can only attend to positions 1..i",
        "• Autoregressive: P(x_i | x_1, ..., x_{i-1}) computed sequentially",
        "• KV cache grows: O(sequence_length × hidden_dim) memory",
        "• Each forward pass requires ALL previous hidden states",
        "• No mathematical way to parallelize across output positions"
    ]
    
    for point in technical_points:
        p = content_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = ""
    
    p = content_frame.add_paragraph()
    p.text = "This is why LLMs cannot be interrupted mid-generation: breaking the causal chain invalidates all subsequent computations."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.italic = True
    
    # Save the presentation
    filename = "slides.pptx"
    prs.save(filename)
    
    
if __name__ == "__main__":
    create_agent1_slides()